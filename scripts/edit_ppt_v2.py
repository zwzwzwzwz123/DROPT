#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
就地改版 组会 PPT -> 20260325ZW组会_v2.pptx (原文件不动)。
- 文字改动: 修 4 硬伤 + 叙事对齐 (保留原字体大小/粗细/颜色)
- 结果图: S16/S20 换成 paper_figures_v2 新图
- 新增页: 三建筑故事 (fig1/2/4/7) 插在实验段
承 docs/PPT_REVISION_GUIDE.md。
"""
import sys, copy
from pptx import Presentation
from pptx.util import Inches, Pt
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

SRC = "docs/20260325ZW组会.pptx"
DST = "docs/20260325ZW组会_v2.pptx"
FIG = "paper_figures_v2"
prs = Presentation(SRC)

def find(slide, shape_id):
    """按 shape_id 递归找 shape (含 group 内)。"""
    def walk(shapes):
        for sh in shapes:
            if sh.shape_id == shape_id:
                return sh
            if sh.shape_type == 6:
                r = walk(sh.shapes)
                if r: return r
        return None
    return walk(slide.shapes)

def set_text(slide, shape_id, lines):
    """替换 shape 文字为 lines(str 或 list[str]); 保留首 run 字体属性。"""
    sh = find(slide, shape_id)
    if sh is None or not sh.has_text_frame:
        print(f"    !! shape {shape_id} 未找到/无文本框"); return
    if isinstance(lines, str): lines = [lines]
    tf = sh.text_frame
    p0 = tf.paragraphs[0]
    donor = p0.runs[0] if p0.runs else None
    f = {}
    if donor is not None:
        f["size"] = donor.font.size; f["bold"] = donor.font.bold; f["name"] = donor.font.name
        try: f["color"] = donor.font.color.rgb
        except Exception: f["color"] = None
    tf.clear()
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run(); run.text = ln
        if donor is not None:
            if f["size"]: run.font.size = f["size"]
            if f["bold"] is not None: run.font.bold = f["bold"]
            if f["name"]: run.font.name = f["name"]
            if f.get("color"):
                try: run.font.color.rgb = f["color"]
                except Exception: pass
    print(f"    ~ shape {shape_id} 文字已改")

def swap_picture(slide, shape_id, img):
    """在原图槽位内换新图, 保持新图宽高比(fit-contain), 居中。"""
    from PIL import Image
    sh = find(slide, shape_id)
    if sh is None: print(f"    !! pic {shape_id} 未找到"); return
    l, t, w, h = sh.left, sh.top, sh.width, sh.height
    iw, ih = Image.open(img).size
    box_ar, img_ar = w / h, iw / ih
    if img_ar > box_ar:            # 图更宽 -> 以宽为准
        nw = w; nh = int(w / img_ar)
    else:                          # 图更高 -> 以高为准
        nh = h; nw = int(h * img_ar)
    nl = l + (w - nw) // 2; nt = t + (h - nh) // 2   # 居中
    sh._element.getparent().remove(sh._element)
    slide.shapes.add_picture(img, nl, nt, nw, nh)
    print(f"    ~ pic {shape_id} 已换图(保比) {img.split('/')[-1]}")

print("开始改版:")
S = prs.slides

# ---- S1 标题: 去掉 Critic-Guided 领衔, 换成 FNO骨干+跨建筑 ----
print("  S1 标题")
set_text(S[0], 8, "Scale-Robust HVAC Control across Buildings via "
                  "Fourier Neural Operator Diffusion Policies")

# ---- S3 Background: 数据中心 -> 商业建筑 HVAC ----
print("  S3 Background")
set_text(S[2], 54, "Buildings account for ~30% of global final energy use, and HVAC is "
                   "the largest single end-use within commercial buildings.")
set_text(S[2], 27, ["Commercial buildings (offices, schools) contain many thermally-coupled "
                    "zones whose HVAC must be coordinated jointly for energy efficiency and comfort.",
                    "We study three BEAR building archetypes spanning 6 / 18 / 25 zones, "
                    "evaluating controllers under a unified protocol across scales.",
                    "Physics-based simulation (BEAR) enables fair, reproducible benchmarking "
                    "of RL controllers without costly co-simulation."])

# ---- S4 传统RL局限: 去掉 6-room 具体化, 弱化多模态 ----
print("  S4 局限")
set_text(S[3], 16, "Struggles to maintain precise coordination across many thermally-coupled "
                   "zones (from 6-zone offices to 25-zone schools) with shared thermal dynamics.")

# ---- S7 Synergy: 多模态弱化, guidance 降级为已验证组件 ----
print("  S7 Synergy")
set_text(S[6], 6, "Diffusion-based Policy: generates smooth, temporally-consistent action "
                  "sequences via an iterative denoising process, avoiding the one-step "
                  "jitter of unimodal Gaussian policies.")
set_text(S[6], 16, "Critic Guidance (validated component, not a contribution): steers denoising "
                   "toward high-reward / safe actions using value gradients; we adopt the "
                   "point-estimate form and quantify its effect via ablation.")

# ---- S12 FNO: 谱截断 通用先验 -> 规模依赖 ----
print("  S12 FNO")
set_text(S[11], 9, ["Transforms the multi-zone action vector to the frequency domain via FFT.",
                    "Keeps low-order spectral modes (a structural prior); truncation is inactive "
                    "on small buildings and only bites as zone count grows.",
                    "Reconstructs the action via inverse FFT — a global, spectrally-consistent mixing."])

# ---- S14 引导: 加诚实 caveat ----
print("  S14 引导")
set_text(S[13], 14, ["Effect 2: Safety & Compliance",
                     "The comfort penalty in the Q-function repels actions that violate the "
                     "±1 C comfort band. Note: this point-estimate guidance is a known biased "
                     "approximation (cf. QGPO) — used here as an empirical component."])

# ---- S16 宏观性能: 换图 (fig1 能耗 + fig2 洼地曲线) + 改文字 ----
print("  S16 宏观性能")
swap_picture(S[15], 12, f"{FIG}/fig1_three_building_energy.png")   # 左大图
swap_picture(S[15], 11, f"{FIG}/fig2_saving_curve_valley.png")     # 右图
set_text(S[15], 6, "Across three building scales (6/18/25 zones), Guided-DiffFNO consistently "
                   "matches or beats the Diff-MLP backbone in energy, with the advantage "
                   "reaching 52% on the largest (School).")
set_text(S[15], 7, "The energy-saving margin is scale-dependent and non-monotonic (12% to 2.2% "
                   "to 52%): a valley at OfficeMedium, where inter-zone coupling is strongest.")
set_text(S[15], 10, "SAC / SAC+MPC baselines are being re-run under a fair, bug-fixed protocol; "
                    "results will report pure-online-RL collapse vs expert-anchored drift.")

# ---- S20 消融: 换图 (fig5 + fig8) + 改口 残差非必需 ----
print("  S20 消融")
swap_picture(S[19], 3, f"{FIG}/fig5_ablation_officesmall.png")     # 左
swap_picture(S[19], 4, f"{FIG}/fig8_tradeoff_scatter.png")         # 右
set_text(S[19], 5, "Ablation (OfficeSmall, 3-seed): removing Critic Guidance markedly raises "
                   "comfort violations; the MLP backbone is clearly worse on both axes.")
set_text(S[19], 12, ["The residual connection is NOT essential for energy: w/o-Residual sits on "
                     "the Pareto front with overlapping std — it mainly aids convergence speed.",
                     "Guided-DiffFNO (Full) sits at the best corner: lowest energy AND lowest "
                     "violations simultaneously."])

# ---- S22 结论: 重写三支柱 ----
print("  S22 结论")
set_text(S[21], 3, [" Cross-scale robustness: an FNO-on-action-axis diffusion policy beats the "
                    "MLP backbone across 6/18/25-zone buildings, most on the largest.",
                    " Component decoupling: controlled 3-building comparison separates backbone "
                    "gain (-38% on School, no guidance) from guidance gain (-22%).",
                    " Honest characterization: spectral truncation helps scale-dependently; "
                    "residual is non-essential; guidance mainly enforces comfort."])

# ---- 新增图页: fig3/4/7/9 (追加在末尾, 你可在 PPT 里拖到实验段) ----
print("  新增 4 张图页 (追加末尾)")
from PIL import Image as _PILImage
# 选一个空白版式(名字含"空白"或占位符最少的)
def pick_blank_layout():
    best, best_n = None, 999
    for lay in prs.slide_layouts:
        n = len(lay.placeholders)
        if ("空白" in lay.name or "Blank" in lay.name) and n < best_n:
            best, best_n = lay, n
    if best is None:  # 兜底: 占位符最少的版式
        best = min(prs.slide_layouts, key=lambda L: len(L.placeholders))
    return best
blank = pick_blank_layout()
print(f"    用空白版式: {blank.name!r} (placeholders={len(blank.placeholders)})")

NEW = [
    ("Per-zone Comfort Violation Rate", f"{FIG}/fig3_violation_rate.png"),
    ("Decoupling Backbone vs Guidance (School)", f"{FIG}/fig4_backbone_guidance_decoupling.png"),
    ("Comfort: Temperature Deviation vs Band", f"{FIG}/fig7_comfort_mean.png"),
    ("Training Stability: FNO vs MLP (School)", f"{FIG}/fig9_training_stability_school.png"),
]
SLIDE_W = Inches(13.333)
for title, img in NEW:
    sl = prs.slides.add_slide(blank)
    # 删掉版式残留的所有占位符(避免"click to add"空框)
    for ph in list(sl.placeholders):
        ph._element.getparent().remove(ph._element)
    # 标题
    tb = sl.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.7))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True
    # 图: 按比例 fit, 目标高 5.3in, 水平居中
    iw, ih = _PILImage.open(img).size
    H = Inches(5.3); Wpx = int(H * iw / ih)
    L = int((SLIDE_W - Wpx) // 2)
    sl.shapes.add_picture(img, L, Inches(1.2), height=H)
    print(f"    + 新页: {title}")

prs.save(DST)
print(f"\n保存 -> {DST}  (共 {len(prs.slides.__iter__.__self__._sldIdLst)} 页)")


